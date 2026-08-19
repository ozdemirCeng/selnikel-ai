import io
import re
from typing import List
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


class EngineeringPowerPointExporter:
    """Exports markdown engineering reports into a widescreen 16:9 presentation slide deck."""

    @classmethod
    def generate_pptx(cls, markdown_text: str, title: str = "Selnikel Mühendislik Sunumu") -> bytes:
        prs = Presentation()
        # Set 16:9 widescreen dimensions
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_slide_layout = prs.slide_layouts[6]

        # 1. Title Slide (Dark Blue Theme)
        title_slide = prs.slides.add_slide(blank_slide_layout)
        background = title_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(8, 15, 30)

        # Title text box
        tx_box = title_slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(3.0))
        tf = tx_box.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = "SELNİKEL ENERJİ"
        p1.font.bold = True
        p1.font.size = Pt(20)
        p1.font.color.rgb = RGBColor(56, 189, 248)

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.bold = True
        p2.font.size = Pt(36)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.space_before = Pt(12)

        p3 = tf.add_paragraph()
        p3.text = "Otomatik Oluşturulan Mühendislik & Teknik Brifing Sunumu"
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(148, 163, 184)
        p3.space_before = Pt(8)

        # 2. Parse Markdown Sections into Content Slides
        lines = markdown_text.split("\n")
        current_section = "Genel Bakış"
        current_bullets: List[str] = []
        current_table: List[List[str]] = []
        in_table = False

        for line in lines:
            trimmed = line.strip()

            if trimmed.startswith("|") and trimmed.endswith("|"):
                if re.match(r"^\|[\s\-:|]+\|$", trimmed):
                    continue
                cols = [c.strip() for c in trimmed.strip("|").split("|")]
                current_table.append(cols)
                in_table = True
                continue
            elif in_table and current_table:
                # Add table slide
                cls._add_table_slide(prs, current_section, current_table)
                current_table = []
                in_table = False

            if trimmed.startswith("## ") or trimmed.startswith("# "):
                # Flush existing bullets
                if current_bullets:
                    cls._add_content_slide(prs, current_section, current_bullets)
                    current_bullets = []
                current_section = trimmed.lstrip("#").strip()
            elif trimmed.startswith("- ") or trimmed.startswith("* "):
                current_bullets.append(trimmed[2:])
            elif trimmed and not trimmed.startswith("---") and not trimmed.startswith("#"):
                current_bullets.append(trimmed)

        # Flush trailing
        if current_bullets:
            cls._add_content_slide(prs, current_section, current_bullets)
        if current_table:
            cls._add_table_slide(prs, current_section, current_table)

        buffer = io.BytesIO()
        prs.save(buffer)
        pptx_bytes = buffer.getvalue()
        buffer.close()
        return pptx_bytes

    @classmethod
    def _add_content_slide(cls, prs: Presentation, title: str, bullets: List[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)

        # Header bar
        header_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        tf_h = header_box.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = title
        p_h.font.bold = True
        p_h.font.size = Pt(24)
        p_h.font.color.rgb = RGBColor(3, 105, 161)

        # Bullets box
        body_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
        tf_b = body_box.text_frame
        tf_b.word_wrap = True

        for idx, bullet in enumerate(bullets[:8]):
            p = tf_b.paragraphs[0] if idx == 0 else tf_b.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(30, 41, 59)
            p.space_after = Pt(8)

    @classmethod
    def _add_table_slide(cls, prs: Presentation, title: str, rows: List[List[str]]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)

        # Header bar
        header_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        tf_h = header_box.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = f"{title} — Parametre Tablosu"
        p_h.font.bold = True
        p_h.font.size = Pt(24)
        p_h.font.color.rgb = RGBColor(3, 105, 161)

        num_rows = min(len(rows), 10)
        num_cols = max(len(r) for r in rows[:num_rows])

        table_shape = slide.shapes.add_table(
            rows=num_rows,
            cols=num_cols,
            left=Inches(1.0),
            top=Inches(2.0),
            width=Inches(11.333),
            height=Inches(4.5),
        )
        table = table_shape.table

        for r_idx in range(num_rows):
            row_data = rows[r_idx]
            for c_idx in range(num_cols):
                val = row_data[c_idx] if c_idx < len(row_data) else ""
                cell = table.cell(r_idx, c_idx)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    if r_idx == 0:
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(255, 255, 255)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(3, 105, 161)
                    else:
                        p.font.color.rgb = RGBColor(30, 41, 59)
                        if r_idx % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
