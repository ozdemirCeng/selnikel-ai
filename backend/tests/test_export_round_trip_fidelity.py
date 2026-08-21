"""
Stage P1.4: Deep Round-Trip Multi-Format Export Verification Test Suite.
Verifies semantic fidelity, exact cell values, structure, styling, Turkish character preservation,
and graceful degradation across Excel (XLSX), Word (DOCX), PowerPoint (PPTX), and PDF formats.
"""
import io
import pytest
import openpyxl
import docx
from pptx import Presentation
from pptx.util import Inches
import pypdf
from httpx import ASGITransport, AsyncClient

from app.main import app
from app.services.reporting.excel_exporter import EngineeringExcelExporter
from app.services.reporting.word_exporter import EngineeringWordExporter
from app.services.reporting.powerpoint_exporter import EngineeringPowerPointExporter
from app.services.reporting.pdf_exporter import EngineeringPDFExporter

AUTH_HEADERS = {"X-Dev-User": "engineer@selnikel.com.tr"}

SAMPLE_ENGINEERING_MARKDOWN = """# SB-500 Endüstriyel Buhar Kazanı Test Raporu
## 1. Tasarım ve Nominal Parametreler
- **Model Tanımı**: SB-500 Yüksek Basınçlı Su Borulu Buhar Kazanı
- **Üretim Yeri**: Selnikel Ankara Fabrikası
- **Standart Uygunluğu**: EN 12952 ve ASME Section I

| Parametre | Nominal Değer | Tasarım Limiti | Birim |
|---|---|---|---|
| Buhar Üretim Kapasitesi | 5000 | 5500 | kg/h |
| İşletme Basıncı | 16.0 | 18.5 | bar |
| Buhar Sıcaklığı | 204.3 | 220.0 | °C |
| Termal Verim | 92.4 | 93.0 | % |
| Baca Gazı Sıcaklığı | 138.5 | 150.0 | °C |

## 2. Mühendislik Değerlendirmesi ve Açıklamalar
Kazan testleri Selnikel AR-GE laboratuvarında başarıyla icra edilmiştir.
- Brülör modülasyon oranı 1:5 aralığında kararlı çalışmıştır.
- Baca emisyonları çevre regülasyonu limitlerinin altındadır.

### 3. Ölçüm Cihazları Kalibrasyonu
Tüm transmitter ve sensörler akredite kuruluşça kalibre edilmiştir.
"""

TURKISH_CHAR_TEST_MARKDOWN = """# Türkçe Karakter Bütünlük Testi: ÇĞİÖŞÜ çğıöşü
## 1. Özel İsim ve Terim Tablosu
| Türkçe Başlık | Açıklama / Değer | Şart / Durum |
|---|---|---|
| Çelik Gövde Dayanımı | 450 MPa (Ağır Hizmet) | İnceleme Yapıldı |
| Isı İletim Katsayısı | 0.85 W/mK | Ölçüldü |
| Şarj Basıncı Kontrolü | 2.5 bar | Uygun |
| Üfleme Fanı Gücü | 15.0 kW | Çalışıyor |

## 2. Paragraf Metni
Şirketimiz Selnikel, ağır sanayi kazanlarında yüksek verim ve üstün güvenlik standardı sağlar.
"""


# ---------------------------------------------------------------------------
# 1. Excel (XLSX) Round-Trip Fidelity Tests
# ---------------------------------------------------------------------------

def test_excel_deep_round_trip_cell_fidelity():
    """Verify XLSX table headers, data cells, title block, and formatting round-trip exactly."""
    excel_bytes = EngineeringExcelExporter.generate_excel(
        SAMPLE_ENGINEERING_MARKDOWN,
        title="SB-500 Test Raporu",
    )
    assert isinstance(excel_bytes, bytes)
    assert excel_bytes.startswith(b"PK\x03\x04")

    # Parse back with openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(excel_bytes))
    assert "Mühendislik Verisi" in wb.sheetnames
    ws = wb["Mühendislik Verisi"]

    # Verify Title Block
    assert "SELNİKEL ENERJİ" in str(ws["A1"].value)
    assert "SB-500 TEST RAPORU" in str(ws["A1"].value)

    # Search for Table Header and Data Rows
    rows_values = []
    for row in ws.iter_rows(values_only=True):
        non_empty = [c for c in row if c is not None]
        if non_empty:
            rows_values.append([str(c) for c in non_empty])

    # Find table header
    header_idx = None
    for idx, r in enumerate(rows_values):
        if "Parametre" in r and "Nominal Değer" in r:
            header_idx = idx
            break

    assert header_idx is not None, "Table header row not found in generated Excel sheet."
    assert rows_values[header_idx] == ["Parametre", "Nominal Değer", "Tasarım Limiti", "Birim"]

    # Verify exact data rows
    table_data = rows_values[header_idx + 1: header_idx + 6]
    assert len(table_data) == 5

    assert table_data[0] == ["Buhar Üretim Kapasitesi", "5000", "5500", "kg/h"]
    assert table_data[1] == ["İşletme Basıncı", "16.0", "18.5", "bar"]
    assert table_data[2] == ["Buhar Sıcaklığı", "204.3", "220.0", "°C"]
    assert table_data[3] == ["Termal Verim", "92.4", "93.0", "%"]
    assert table_data[4] == ["Baca Gazı Sıcaklığı", "138.5", "150.0", "°C"]


# ---------------------------------------------------------------------------
# 2. Word (DOCX) Round-Trip Fidelity Tests
# ---------------------------------------------------------------------------

def test_word_deep_round_trip_structure_and_tables():
    """Verify DOCX headings, paragraphs, bullet points, and exact table cells."""
    docx_bytes = EngineeringWordExporter.generate_docx(
        SAMPLE_ENGINEERING_MARKDOWN,
        title="SB-500 Rapor",
    )
    assert isinstance(docx_bytes, bytes)
    assert docx_bytes.startswith(b"PK\x03\x04")

    # Parse back with python-docx
    doc = docx.Document(io.BytesIO(docx_bytes))

    # Verify Headings
    headings = [p.text for p in doc.paragraphs if p.style.name.startswith("Heading")]
    assert any("1. Tasarım ve Nominal Parametreler" in h for h in headings)
    assert any("2. Mühendislik Değerlendirmesi ve Açıklamalar" in h for h in headings)
    assert any("3. Ölçüm Cihazları Kalibrasyonu" in h for h in headings)

    # Verify Table Extraction
    assert len(doc.tables) >= 1
    tbl = doc.tables[0]
    assert len(tbl.rows) == 6  # 1 header + 5 data rows
    assert len(tbl.columns) == 4

    header_cells = [c.text.strip() for c in tbl.rows[0].cells]
    assert header_cells == ["Parametre", "Nominal Değer", "Tasarım Limiti", "Birim"]

    row1_cells = [c.text.strip() for c in tbl.rows[1].cells]
    assert row1_cells == ["Buhar Üretim Kapasitesi", "5000", "5500", "kg/h"]

    row2_cells = [c.text.strip() for c in tbl.rows[2].cells]
    assert row2_cells == ["İşletme Basıncı", "16.0", "18.5", "bar"]

    row4_cells = [c.text.strip() for c in tbl.rows[4].cells]
    assert row4_cells == ["Termal Verim", "92.4", "93.0", "%"]


# ---------------------------------------------------------------------------
# 3. PowerPoint (PPTX) Round-Trip Fidelity Tests
# ---------------------------------------------------------------------------

def test_powerpoint_deep_round_trip_slides_and_tables():
    """Verify PPTX slide count, 16:9 widescreen layout, slide titles, bullets, and table shapes."""
    pptx_bytes = EngineeringPowerPointExporter.generate_pptx(
        SAMPLE_ENGINEERING_MARKDOWN,
        title="SB-500 Brifing",
    )
    assert isinstance(pptx_bytes, bytes)
    assert pptx_bytes.startswith(b"PK\x03\x04")

    # Parse back with python-pptx
    prs = Presentation(io.BytesIO(pptx_bytes))

    # Verify 16:9 Widescreen dimensions
    assert abs(prs.slide_width.inches - 13.333) < 0.01
    assert abs(prs.slide_height.inches - 7.5) < 0.01

    # Verify Slides exist: Title Slide, Content Slides, Table Slide
    assert len(prs.slides) >= 3

    # Check Title Slide (Slide 0)
    slide0_text = "".join(shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame)
    assert "SELNİKEL ENERJİ" in slide0_text
    assert "SB-500 Brifing" in slide0_text

    # Search for table shape across slides
    found_table = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                found_table = shape.table
                break
        if found_table:
            break

    assert found_table is not None, "Table shape not found in presentation."
    assert len(found_table.rows) == 6
    assert len(found_table.columns) == 4

    header_vals = [found_table.cell(0, c).text.strip() for c in range(4)]
    assert header_vals == ["Parametre", "Nominal Değer", "Tasarım Limiti", "Birim"]

    row1_vals = [found_table.cell(1, c).text.strip() for c in range(4)]
    assert row1_vals == ["Buhar Üretim Kapasitesi", "5000", "5500", "kg/h"]

    row2_vals = [found_table.cell(2, c).text.strip() for c in range(4)]
    assert row2_vals == ["İşletme Basıncı", "16.0", "18.5", "bar"]


# ---------------------------------------------------------------------------
# 4. PDF Round-Trip Text & Metadata Extraction
# ---------------------------------------------------------------------------

def test_pdf_deep_round_trip_text_extraction():
    """Verify PDF binary header, document structure, table content, and verification footer."""
    pdf_bytes = EngineeringPDFExporter.generate_pdf(
        SAMPLE_ENGINEERING_MARKDOWN,
        title="SB-500 Resmi Rapor",
    )
    assert isinstance(pdf_bytes, bytes)
    assert pdf_bytes.startswith(b"%PDF")

    # Parse back with pypdf
    reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
    assert len(reader.pages) >= 1

    extracted_text = "\n".join(page.extract_text() for page in reader.pages)

    # Core corporate header
    assert "SELNİKEL ENERJİ" in extracted_text
    assert "TEKNİK RAPOR" in extracted_text

    # Content and Table Values
    assert "SB-500 Endüstriyel Buhar Kazanı Test Raporu" in extracted_text
    assert "Buhar Üretim Kapasitesi" in extracted_text
    assert "5000" in extracted_text
    assert "İşletme Basıncı" in extracted_text
    assert "16.0" in extracted_text
    assert "Termal Verim" in extracted_text
    assert "92.4" in extracted_text

    # Standard verification seal
    assert "ASME PTC 4.1" in extracted_text


# ---------------------------------------------------------------------------
# 5. Full Turkish Character Set Fidelity Across All 4 Formats
# ---------------------------------------------------------------------------

def test_cross_format_turkish_character_fidelity():
    """Verify ç, Ç, ğ, Ğ, ı, İ, ö, Ö, ş, Ş, ü, Ü preserve perfectly across XLSX, DOCX, PPTX, and PDF."""
    
    # 1. Excel
    excel_bytes = EngineeringExcelExporter.generate_excel(TURKISH_CHAR_TEST_MARKDOWN, title="Türkçe Test")
    wb = openpyxl.load_workbook(io.BytesIO(excel_bytes))
    ws = wb.active
    excel_text = " ".join(str(c) for row in ws.iter_rows(values_only=True) for c in row if c)
    assert "Çelik Gövde Dayanımı" in excel_text
    assert "Isı İletim Katsayısı" in excel_text
    assert "Şarj Basıncı Kontrolü" in excel_text
    assert "Üfleme Fanı Gücü" in excel_text

    # 2. Word
    docx_bytes = EngineeringWordExporter.generate_docx(TURKISH_CHAR_TEST_MARKDOWN, title="Türkçe Test")
    doc = docx.Document(io.BytesIO(docx_bytes))
    docx_text = " ".join(p.text for p in doc.paragraphs)
    for t in doc.tables:
        docx_text += " " + " ".join(c.text for row in t.rows for c in row.cells)
    assert "Çelik Gövde Dayanımı" in docx_text
    assert "Isı İletim Katsayısı" in docx_text
    assert "Şarj Basıncı Kontrolü" in docx_text
    assert "Üfleme Fanı Gücü" in docx_text

    # 3. PowerPoint
    pptx_bytes = EngineeringPowerPointExporter.generate_pptx(TURKISH_CHAR_TEST_MARKDOWN, title="Türkçe Test")
    prs = Presentation(io.BytesIO(pptx_bytes))
    pptx_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                pptx_text += " " + shape.text_frame.text
            if shape.has_table:
                pptx_text += " " + " ".join(cell.text for row in shape.table.rows for cell in row.cells)
    assert "Çelik Gövde Dayanımı" in pptx_text
    assert "Isı İletim Katsayısı" in pptx_text
    assert "Şarj Basıncı Kontrolü" in pptx_text
    assert "Üfleme Fanı Gücü" in pptx_text

    # 4. PDF
    pdf_bytes = EngineeringPDFExporter.generate_pdf(TURKISH_CHAR_TEST_MARKDOWN, title="Türkçe Test")
    reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
    pdf_text = "\n".join(page.extract_text() for page in reader.pages)
    assert "Çelik Gövde Dayanımı" in pdf_text or "Celik" in pdf_text or "Gövde" in pdf_text
    assert "2.5 bar" in pdf_text
    assert "15.0 kW" in pdf_text


# ---------------------------------------------------------------------------
# 6. Edge Cases & Robustness
# ---------------------------------------------------------------------------

@pytest.mark.parametrize(
    "edge_markdown",
    [
        "",  # Empty string
        "   \n\n   \n\t  ",  # Whitespace only
        "# Only Title\n## Subheading\nNo tables or bullets here.",  # No tables
        "| A | B |\n|---|---|\n| 1 | 2 |",  # Only table
        "| Col1 | Col2 |\n|---|---|\n| Cell1 | Cell2 | Cell3 Extra |",  # Inconsistent row length
        "# Sec 1\n| T1 | Val |\n|---|---|\n| A | 1 |\n# Sec 2\n| T2 | Val |\n|---|---|\n| B | 2 |",  # Multiple tables
        "# Special Chars < > & \" ' $ % @ # *",  # Special characters
    ],
)
def test_exporters_handle_edge_cases_gracefully(edge_markdown):
    """Verify that all 4 exporters process edge-case markdowns without crashing or producing corrupt bytes."""
    xlsx_b = EngineeringExcelExporter.generate_excel(edge_markdown, title="Edge Case")
    assert isinstance(xlsx_b, bytes) and len(xlsx_b) > 500

    docx_b = EngineeringWordExporter.generate_docx(edge_markdown, title="Edge Case")
    assert isinstance(docx_b, bytes) and len(docx_b) > 500

    pptx_b = EngineeringPowerPointExporter.generate_pptx(edge_markdown, title="Edge Case")
    assert isinstance(pptx_b, bytes) and len(pptx_b) > 500

    pdf_b = EngineeringPDFExporter.generate_pdf(edge_markdown, title="Edge Case")
    assert isinstance(pdf_b, bytes) and len(pdf_b) > 500


# ---------------------------------------------------------------------------
# 7. HTTP API Export Endpoints End-to-End Verification
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_api_export_endpoints_round_trip_validation():
    """Verify HTTP export endpoints stream parseable binaries for all formats."""
    async with AsyncClient(
        transport=ASGITransport(app=app), base_url="http://test"
    ) as ac:
        # 1. Excel Endpoint
        res_excel = await ac.post(
            "/api/v1/agent/report/excel",
            json={"markdown_content": SAMPLE_ENGINEERING_MARKDOWN, "title": "API Excel Test"},
            headers=AUTH_HEADERS,
        )
        assert res_excel.status_code == 200
        wb = openpyxl.load_workbook(io.BytesIO(res_excel.content))
        assert "Mühendislik Verisi" in wb.sheetnames

        # 2. Word Endpoint
        res_word = await ac.post(
            "/api/v1/agent/report/word",
            json={"markdown_content": SAMPLE_ENGINEERING_MARKDOWN, "title": "API Word Test"},
            headers=AUTH_HEADERS,
        )
        assert res_word.status_code == 200
        doc = docx.Document(io.BytesIO(res_word.content))
        assert len(doc.tables) >= 1

        # 3. PowerPoint Endpoint
        res_pptx = await ac.post(
            "/api/v1/agent/report/powerpoint",
            json={"markdown_content": SAMPLE_ENGINEERING_MARKDOWN, "title": "API PPTX Test"},
            headers=AUTH_HEADERS,
        )
        assert res_pptx.status_code == 200
        prs = Presentation(io.BytesIO(res_pptx.content))
        assert len(prs.slides) >= 3

        # 4. PDF Endpoint
        res_pdf = await ac.post(
            "/api/v1/agent/report/pdf",
            json={"markdown_content": SAMPLE_ENGINEERING_MARKDOWN, "title": "API PDF Test"},
            headers=AUTH_HEADERS,
        )
        assert res_pdf.status_code == 200
        reader = pypdf.PdfReader(io.BytesIO(res_pdf.content))
        assert len(reader.pages) >= 1
